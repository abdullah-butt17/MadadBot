import gradio as gr
import os
from typing import List, Dict, Any
from langchain_groq import ChatGroq
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document
from langchain_core.tools import tool
from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langgraph.prebuilt import create_react_agent
from langgraph.checkpoint.memory import MemorySaver
from pydantic import BaseModel, Field
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import pytesseract
import requests
import re

# Global variables
vectorstore = None
agent_executor = None
memory = MemorySaver()
user_name = None
embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

# Groq client initialization
def get_groq_client():
    from google.colab import userdata
    api_key = userdata.get('MadadBot')
    if not api_key:
        raise ValueError("Groq API key not found")

    return ChatGroq(
        api_key=api_key,
        model="llama-3.3-70b-versatile",
        temperature=0.3,
        max_tokens=2000,
    )

# Document extraction functions
def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = pypdf.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(file_path):
    doc = DocxDocument(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_image(file_path):
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text
    except:
        return ""

def process_uploaded_files(files):
    if not files:
        return ""
    all_text = ""
    for file in files:
        file_path = file.name
        file_extension = file_path.split('.')[-1].lower()
        try:
            if file_extension == 'pdf':
                all_text += extract_text_from_pdf(file_path) + "\n\n"
            elif file_extension == 'docx':
                all_text += extract_text_from_docx(file_path) + "\n\n"
            elif file_extension == 'txt':
                all_text += extract_text_from_txt(file_path) + "\n\n"
            elif file_extension == 'pptx':
                all_text += extract_text_from_pptx(file_path) + "\n\n"
            elif file_extension in ['png', 'jpg', 'jpeg']:
                all_text += extract_text_from_image(file_path) + "\n\n"
        except Exception:
            continue
    return all_text

def create_vectorstore(text):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    documents = [Document(page_content=chunk) for chunk in chunks]
    vectorstore = FAISS.from_documents(documents, embeddings)
    return vectorstore

# Pydantic models for tool inputs
class RetrieveDocsInput(BaseModel):
    query: str = Field(description="The search query to find relevant information in uploaded documents")

class SearchWebInput(BaseModel):
    query: str = Field(description="The search query to find information on the internet")

# Tool Definitions
@tool(args_schema=RetrieveDocsInput)
def retrieve_documents(query: str) -> str:
    """Search uploaded documents for relevant information."""
    global vectorstore
    if not vectorstore:
        return "No documents have been uploaded yet."

    try:
        docs = vectorstore.similarity_search(query, k=3)
        if not docs:
            return "No relevant information found in uploaded documents."

        context = "\n\n---\n\n".join([doc.page_content for doc in docs])
        return f"Found relevant information in uploaded documents:\n\n{context}"
    except Exception as e:
        return f"Error searching documents: {str(e)}"

@tool(args_schema=SearchWebInput)
def search_web(query: str) -> str:
    """Search the internet for current information."""

    def search_wikipedia(q):
        try:
            search_url = "https://en.wikipedia.org/w/api.php"
            search_params = {
                "action": "query",
                "list": "search",
                "srsearch": q,
                "format": "json",
                "srlimit": 2
            }
            response = requests.get(search_url, params=search_params, timeout=10)
            search_data = response.json()
            results = []
            if 'query' in search_data and 'search' in search_data['query']:
                for item in search_data['query']['search'][:2]:
                    page_title = item['title']
                    content_params = {
                        "action": "query",
                        "prop": "extracts",
                        "exintro": True,
                        "explaintext": True,
                        "titles": page_title,
                        "format": "json"
                    }
                    content_response = requests.get(search_url, params=content_params, timeout=10)
                    content_data = content_response.json()
                    pages = content_data.get('query', {}).get('pages', {})
                    for page_id, page_data in pages.items():
                        if 'extract' in page_data:
                            results.append({
                                'url': f"https://en.wikipedia.org/wiki/{page_title.replace(' ', '_')}",
                                'title': page_title,
                                'text': page_data['extract'][:1500]
                            })
            return results
        except:
            return []

    def search_duckduckgo(q):
        try:
            url = "https://api.duckduckgo.com/"
            params = {"q": q, "format": "json", "no_html": 1, "skip_disambig": 1}
            headers = {'User-Agent': 'Mozilla/5.0'}
            response = requests.get(url, params=params, headers=headers, timeout=10)
            data = response.json()
            results = []
            if data.get('Abstract'):
                results.append({
                    'url': data.get('AbstractURL', 'https://duckduckgo.com'),
                    'title': data.get('Heading', 'DuckDuckGo Result'),
                    'text': data.get('Abstract', '')
                })
            return results
        except:
            return []

    all_results = []
    ddg_results = search_duckduckgo(query)
    wiki_results = search_wikipedia(query)
    all_results.extend(ddg_results)
    all_results.extend(wiki_results)

    if not all_results:
        return "No relevant information found on the web."

    formatted_results = "Web search results:\n\n"
    for i, r in enumerate(all_results, 1):
        formatted_results += f"{i}. **{r['title']}**\n"
        formatted_results += f"   URL: {r['url']}\n"
        formatted_results += f"   Content: {r['text'][:500]}...\n\n"

    return formatted_results

# Agent initialization
def initialize_agent():
    global agent_executor
    try:
        llm = get_groq_client()
        tools = [retrieve_documents, search_web]
        agent_executor = create_react_agent(llm, tools, checkpointer=memory)
        return "✅ Agent initialized successfully!"
    except Exception as e:
        return f"❌ Error initializing agent: {str(e)}"

# Process documents
def process_documents(files):
    global vectorstore
    if not files:
        return "❌ Please upload at least one document."
    try:
        all_text = process_uploaded_files(files)
        if all_text.strip():
            vectorstore = create_vectorstore(all_text)
            return f"✅ Successfully processed {len(files)} document(s)! Vector store created with {len(all_text)} characters."
        else:
            return "❌ No text could be extracted from the documents."
    except Exception as e:
        return f"❌ Error processing documents: {str(e)}"

# Answer question function
def answer_question(question, history):
    global agent_executor, user_name

    # Handle name recognition
    match = re.search(r"my name is (\w+)", question, re.I)
    if match:
        user_name = match.group(1)
        response_text = f"Nice to meet you, {user_name}! 😊"
        history.append({"role": "assistant", "content": response_text})
        return history, ""

    if re.search(r"\bmy name\b", question.lower()) or re.search(r"\bremember\b", question.lower()):
        if user_name:
            response_text = f"Yes, I remember your name is {user_name}! 👋"
        else:
            response_text = "I'm sorry, I don't recall your name yet. Could you please tell me again?"
        history.append({"role": "assistant", "content": response_text})
        return history, ""

    if not agent_executor:
        init_msg = initialize_agent()
        if "Error" in init_msg:
            return history + [{"role": "assistant", "content": init_msg}], ""

    if not question.strip():
        return history + [{"role": "assistant", "content": "⚠️ Please enter a question."}], ""

    try:
        history.append({"role": "user", "content": question})
        thread_id = "madadbot_conversation"
        config = {"configurable": {"thread_id": thread_id}}

        system_msg = SystemMessage(content="""You are MadadBot, an intelligent AI assistant.

Your process:
1. Analyze questions carefully and check conversation history
2. For follow-up questions, use previous context
3. Use retrieve_documents for document questions
4. Use search_web for current events or general knowledge
5. Provide clear, well-structured answers with markdown formatting

Format answers with headers, bullet points, and citations when needed.""")

        recent_msgs = []
        context_window = min(10, len(history) - 1)

        for msg in history[-context_window:]:
            if msg != history[-1]:
                if msg["role"] == "user":
                    recent_msgs.append(HumanMessage(content=msg["content"]))
                elif msg["role"] == "assistant":
                    recent_msgs.append(AIMessage(content=msg["content"][:1000]))

        recent_msgs.append(HumanMessage(content=question))

        try:
            response = agent_executor.invoke(
                {"messages": [system_msg] + recent_msgs},
                config=config
            )
            final_message = response["messages"][-1].content

        except Exception as tool_error:
            print(f"Agent error: {tool_error}. Using fallback.")

            follow_up_keywords = ["it", "that", "this", "more about", "tell me more"]
            is_follow_up = any(keyword in question.lower() for keyword in follow_up_keywords)

            context = ""

            if is_follow_up and len(history) > 2:
                context += "\n\n=== Previous Context ===\n"
                for msg in history[-6:-1]:
                    if msg["role"] == "user":
                        context += f"User: {msg['content']}\n"
                    elif msg["role"] == "assistant":
                        context += f"Assistant: {msg['content'][:800]}\n\n"

            if vectorstore:
                doc_results = retrieve_documents.invoke({"query": question})
                if "No relevant information" not in doc_results:
                    context += f"\n\nDocuments:\n{doc_results}\n"

            web_results = search_web.invoke({"query": question})
            context += f"\n\nWeb:\n{web_results}\n"

            llm = get_groq_client()
            prompt = f"""Based on context and history, answer the question comprehensively.

Context:
{context}

Question: {question}

Provide a clear, well-structured answer with markdown formatting."""

            llm_response = llm.invoke([HumanMessage(content=prompt)])
            final_message = llm_response.content

        formatted_response = f"""{final_message}

---
*Generated by MadadBot*"""

        history.append({"role": "assistant", "content": formatted_response})
        return history, ""

    except Exception as e:
        error_msg = f"❌ Error: {str(e)}\n\nPlease try again or check your API key."
        history.append({"role": "assistant", "content": error_msg})
        return history, ""

# ChatGPT-style CSS
custom_css = """
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');

:root {
    --bg-primary: #0c0e1a;
    --bg-secondary: #1a1d2e;
    --bg-tertiary: #252836;
    --text-primary: #ececf1;
    --text-secondary: #8e8ea0;
    --border-color: #2f3241;
    --accent-color: #10a37f;
    --hover-bg: #2a2b38;
}

body, .gradio-container {
    font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif !important;
    background: var(--bg-primary) !important;
    color: var(--text-primary) !important;
}

/* Hide Gradio branding */
footer {
    display: none !important;
}

/* Header styling */
.header-container {
    background: var(--bg-secondary);
    border-bottom: 1px solid var(--border-color);
    padding: 16px 24px;
    margin: -8px -8px 20px -8px;
}

.header-title {
    font-size: 24px;
    font-weight: 600;
    background: linear-gradient(90deg, #10a37f 0%, #1a7f64 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    margin-bottom: 4px;
}

.header-subtitle {
    font-size: 14px;
    color: var(--text-secondary);
}

/* Main layout */
.gr-row {
    gap: 20px !important;
}

.gr-column {
    background: var(--bg-secondary) !important;
    border-radius: 12px !important;
    padding: 20px !important;
    border: 1px solid var(--border-color) !important;
}

/* Labels */
label {
    color: var(--text-primary) !important;
    font-weight: 500 !important;
    font-size: 14px !important;
    margin-bottom: 8px !important;
}

/* File upload */
.gr-file {
    background: var(--bg-tertiary) !important;
    border: 2px dashed var(--border-color) !important;
    border-radius: 10px !important;
    padding: 20px !important;
    transition: all 0.2s !important;
}

.gr-file:hover {
    border-color: var(--accent-color) !important;
    background: rgba(16, 163, 127, 0.05) !important;
}

/* Buttons */
.gr-button {
    border-radius: 8px !important;
    font-weight: 500 !important;
    font-size: 14px !important;
    padding: 10px 16px !important;
    transition: all 0.2s !important;
    border: none !important;
}

.gr-button-primary {
    background: var(--accent-color) !important;
    color: white !important;
}

.gr-button-primary:hover {
    background: #0d8c6f !important;
    transform: translateY(-1px);
    box-shadow: 0 4px 12px rgba(16, 163, 127, 0.3) !important;
}

.gr-button-secondary {
    background: var(--bg-tertiary) !important;
    color: var(--text-primary) !important;
    border: 1px solid var(--border-color) !important;
}

.gr-button-secondary:hover {
    background: var(--hover-bg) !important;
}

/* Textbox */
.gr-textbox {
    background: var(--bg-tertiary) !important;
    border: 1px solid var(--border-color) !important;
    border-radius: 8px !important;
    color: var(--text-primary) !important;
    font-size: 14px !important;
}

.gr-textbox:focus {
    border-color: var(--accent-color) !important;
    box-shadow: 0 0 0 2px rgba(16, 163, 127, 0.1) !important;
}

/* Chatbot */
.gr-chatbot {
    background: var(--bg-primary) !important;
    border: 1px solid var(--border-color) !important;
    border-radius: 12px !important;
    min-height: 500px !important;
}

.gr-chatbot .message {
    padding: 16px !important;
    margin: 8px 0 !important;
    border-radius: 10px !important;
    font-size: 15px !important;
    line-height: 1.6 !important;
}

.gr-chatbot .message.user {
    background: var(--bg-tertiary) !important;
    color: var(--text-primary) !important;
    margin-left: 48px !important;
    border: 1px solid var(--border-color) !important;
}

.gr-chatbot .message.assistant {
    background: var(--bg-secondary) !important;
    color: var(--text-primary) !important;
    margin-right: 48px !important;
    border: 1px solid var(--border-color) !important;
}

/* Markdown in messages */
.gr-chatbot .message h1,
.gr-chatbot .message h2,
.gr-chatbot .message h3 {
    color: var(--accent-color) !important;
    margin: 12px 0 8px !important;
    font-weight: 600 !important;
}

.gr-chatbot .message ul,
.gr-chatbot .message ol {
    margin: 8px 0 8px 20px !important;
}

.gr-chatbot .message code {
    background: var(--bg-tertiary) !important;
    padding: 2px 6px !important;
    border-radius: 4px !important;
    font-family: 'Courier New', monospace !important;
    font-size: 13px !important;
    color: #e96900 !important;
}

.gr-chatbot .message pre {
    background: var(--bg-tertiary) !important;
    padding: 12px !important;
    border-radius: 8px !important;
    overflow-x: auto !important;
    margin: 12px 0 !important;
    border: 1px solid var(--border-color) !important;
}

/* Status output */
#process_output {
    background: var(--bg-tertiary) !important;
    border: 1px solid var(--border-color) !important;
    border-radius: 8px !important;
    color: var(--text-secondary) !important;
    padding: 12px !important;
    font-size: 13px !important;
    font-family: 'Courier New', monospace !important;
}

/* Scrollbars */
::-webkit-scrollbar {
    width: 8px;
    height: 8px;
}

::-webkit-scrollbar-track {
    background: var(--bg-primary);
}

::-webkit-scrollbar-thumb {
    background: var(--border-color);
    border-radius: 4px;
}

::-webkit-scrollbar-thumb:hover {
    background: #4a4b5c;
}

/* Welcome banner */
.welcome-banner {
    background: linear-gradient(135deg, var(--bg-secondary) 0%, var(--bg-tertiary) 100%);
    border: 1px solid var(--border-color);
    border-radius: 12px;
    padding: 24px;
    margin-bottom: 20px;
    text-align: center;
}

.welcome-title {
    font-size: 28px;
    font-weight: 700;
    background: linear-gradient(90deg, #10a37f 0%, #1a7f64 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    margin-bottom: 8px;
}

.welcome-desc {
    color: var(--text-secondary);
    font-size: 15px;
    line-height: 1.5;
}

/* Feature cards */
.feature-grid {
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
    gap: 12px;
    margin-top: 16px;
}

.feature-card {
    background: var(--bg-tertiary);
    border: 1px solid var(--border-color);
    border-radius: 10px;
    padding: 16px;
    text-align: center;
    transition: all 0.2s;
}

.feature-card:hover {
    border-color: var(--accent-color);
    transform: translateY(-2px);
}

.feature-icon {
    font-size: 32px;
    margin-bottom: 8px;
}

.feature-title {
    font-size: 14px;
    font-weight: 600;
    color: var(--text-primary);
    margin-bottom: 4px;
}

.feature-desc {
    font-size: 12px;
    color: var(--text-secondary);
}

/* Responsive */
@media (max-width: 768px) {
    .gr-row {
        flex-direction: column !important;
    }

    .gr-chatbot .message.user {
        margin-left: 0 !important;
    }

    .gr-chatbot .message.assistant {
        margin-right: 0 !important;
    }
}
"""

# Gradio Interface
with gr.Blocks(css=custom_css, theme=gr.themes.Base(), title="MadadBot - AI Assistant") as demo:
    # Header
    gr.HTML("""
        <div class="header-container">
            <div class="header-title">🤖 MadadBot - Agentic RAG System</div>
            <div class="header-subtitle">Powered by LangGraph, Groq AI & Llama 3.3 70B</div>
        </div>
    """)

    # Welcome Banner
    gr.HTML("""
        <div class="welcome-banner">
            <div class="welcome-title">Welcome to MadadBot</div>
            <div class="welcome-desc">
                Your intelligent AI assistant with document processing and web search capabilities.
                Upload documents to analyze or ask questions to get started.
            </div>
            <div class="feature-grid">
                <div class="feature-card">
                    <div class="feature-icon">📄</div>
                    <div class="feature-title">Document Analysis</div>
                    <div class="feature-desc">Process PDFs, DOCX, TXT & more</div>
                </div>
                <div class="feature-card">
                    <div class="feature-icon">🔍</div>
                    <div class="feature-title">Smart Search</div>
                    <div class="feature-desc">Semantic search with RAG</div>
                </div>
                <div class="feature-card">
                    <div class="feature-icon">🌐</div>
                    <div class="feature-title">Web Integration</div>
                    <div class="feature-desc">Real-time web information</div>
                </div>
                <div class="feature-card">
                    <div class="feature-icon">🤖</div>
                    <div class="feature-title">Agentic AI</div>
                    <div class="feature-desc">Autonomous reasoning & tools</div>
                </div>
            </div>
        </div>
    """)

    # Main Interface
    with gr.Row():
        # Left Column - Document Upload
        with gr.Column(scale=1):
            gr.Markdown("### 📚 Document Management")
            file_upload = gr.File(
                label="Upload Documents",
                file_count="multiple",
                file_types=[".pdf", ".docx", ".txt", ".pptx", ".png", ".jpg", ".jpeg"],
                height=150
            )
            process_btn = gr.Button("🚀 Process Documents", variant="primary", size="lg")
            process_output = gr.Textbox(
                label="Processing Status",
                interactive=False,
                lines=4,
                elem_id="process_output",
                placeholder="Upload files and click 'Process Documents' to begin..."
            )

        # Right Column - Chat Interface
        with gr.Column(scale=2):
            gr.Markdown("### 💬 Chat with MadadBot")
            chatbot = gr.Chatbot(
                label="",
                type="messages",
                height=520,
                show_label=False
            )

            with gr.Row():
                msg = gr.Textbox(
                    label="",
                    placeholder="Message MadadBot... (Press Enter to send, Shift+Enter for new line)",
                    lines=2,
                    max_lines=6,
                    show_label=False,
                    scale=4
                )
                with gr.Column(scale=1, min_width=100):
                    submit = gr.Button("🔍 Send", variant="primary", size="lg")
                    clear = gr.Button("🗑️ Clear", variant="secondary", size="sm")

    # Footer info
    gr.HTML("""
        <div style="text-align: center; margin-top: 20px; padding: 16px; background: var(--bg-secondary); border-radius: 8px; border: 1px solid var(--border-color);">
            <p style="color: var(--text-secondary); font-size: 13px; margin: 0;">
                💡 <strong>Tip:</strong> Upload your documents first for better context-aware responses.
                MadadBot can analyze PDFs, Word documents, presentations, and images with text.
            </p>
        </div>
    """)

    # Event handlers
    process_btn.click(process_documents, inputs=[file_upload], outputs=process_output)
    submit.click(answer_question, inputs=[msg, chatbot], outputs=[chatbot, msg])
    msg.submit(answer_question, inputs=[msg, chatbot], outputs=[chatbot, msg])
    clear.click(lambda: [], outputs=chatbot)

# Initialize agent on startup
print("🚀 Initializing MadadBot Agent...")
init_result = initialize_agent()
print(init_result)

if __name__ == "__main__":
    print("\n" + "="*60)
    print("🤖 MadadBot - Agentic RAG System")
    print("="*60)
    print("📦 Starting Gradio interface...")
    demo.launch(share=True, show_error=True)
